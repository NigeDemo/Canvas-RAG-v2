"""Content processing module for PDFs and images."""

import base64
from io import BytesIO
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
import json

from PIL import Image
import fitz  # PyMuPDF
import pypdf
from pptx import Presentation
from bs4 import BeautifulSoup
import re

from ..config.settings import settings
from ..utils.logger import get_logger
from ..vision.vision_processor import VisionProcessor

logger = get_logger(__name__)

class ContentProcessor:
    """Processes various content types for the RAG system."""
    
    def __init__(self, enable_vision: bool = True):
        """Initialize the content processor."""
        self.max_image_size = settings.max_image_size
        self.pdf_dpi = settings.pdf_dpi
        self.chunk_size = settings.chunk_size
        self.chunk_overlap = settings.chunk_overlap
        self.enable_vision = enable_vision
        
        # Initialize vision processor if enabled
        if enable_vision:
            try:
                self.vision_processor = VisionProcessor(
                    primary_provider=settings.vision_primary_provider,
                    fallback_provider=settings.vision_fallback_provider,
                    use_cache=True
                )
                logger.info("Vision AI enabled for content processing")
            except Exception as e:
                logger.warning(f"Failed to initialize vision processor: {e}")
                self.vision_processor = None
                self.enable_vision = False
        else:
            self.vision_processor = None
            logger.info("Vision AI disabled for content processing")
    
    def _analyze_image_content(self, image_input, alt_text: str = "") -> Dict[str, Any]:
        """
        Analyze image content using vision AI.
        
        Args:
            image_input: URL, file path, or base64 data of the image to analyze
            alt_text: Existing alt text for context
            
        Returns:
            Dictionary with vision analysis results
        """
        if not self.enable_vision or not self.vision_processor:
            return {}
        
        try:
            # Create architectural drawing analysis prompt
            prompt = f"""Analyze this architectural drawing/image. Provide:
1. Type of drawing (floor plan, elevation, section, detail, etc.)
2. Key architectural elements visible
3. Scale information if visible
4. Dimensions or measurements shown
5. Construction details and materials
6. Any text or annotations visible

Context from alt text: {alt_text}

Focus on technical architectural content that would be relevant for construction and design education."""
            
            # Analyze the image
            analysis = self.vision_processor.analyze_image(
                image_input=image_input,
                prompt=prompt
            )
            
            if analysis and analysis.get('success'):
                return {
                    'vision_analysis': analysis.get('analysis', ''),
                    'analysis_provider': analysis.get('provider', ''),
                    'analysis_timestamp': analysis.get('timestamp', ''),
                    'drawing_type': self._extract_drawing_type(analysis.get('analysis', '')),
                    'has_vision_analysis': True
                }
            
        except Exception as e:
            logger.warning(f"Vision analysis failed: {e}")
        
        return {'has_vision_analysis': False}
    
    def _extract_drawing_type(self, analysis_text: str) -> str:
        """Extract drawing type from vision analysis."""
        text_lower = analysis_text.lower()
        
        # Check for drawing type keywords
        drawing_types = {
            'floor plan': ['floor plan', 'plan view', 'layout'],
            'elevation': ['elevation', 'facade', 'front view'],
            'section': ['section', 'cross section', 'sectional'],
            'detail': ['detail', 'construction detail', 'close-up'],
            'site plan': ['site plan', 'site layout', 'topographical'],
            'perspective': ['perspective', '3d view', 'isometric'],
            'diagram': ['diagram', 'schematic', 'chart']
        }
        
        for drawing_type, keywords in drawing_types.items():
            if any(keyword in text_lower for keyword in keywords):
                return drawing_type
        
        return 'unknown'
    
    def resize_image(self, image: Image.Image) -> Image.Image:
        """
        Resize image while maintaining aspect ratio.
        
        Args:
            image: PIL Image object
            
        Returns:
            Resized image
        """
        width, height = image.size
        
        if width <= self.max_image_size and height <= self.max_image_size:
            return image
        
        # Calculate new dimensions
        if width > height:
            new_width = self.max_image_size
            new_height = int((height * self.max_image_size) / width)
        else:
            new_height = self.max_image_size
            new_width = int((width * self.max_image_size) / height)
        
        return image.resize((new_width, new_height), Image.Resampling.LANCZOS)
    
    def image_to_base64(self, image: Image.Image) -> str:
        """
        Convert PIL Image to base64 string.
        
        Args:
            image: PIL Image object
            
        Returns:
            Base64 encoded image string
        """
        buffer = BytesIO()
        image.save(buffer, format='PNG')
        image_bytes = buffer.getvalue()
        return base64.b64encode(image_bytes).decode('utf-8')
    
    def process_pdf(self, pdf_path: Path) -> List[Dict[str, Any]]:
        """
        Process PDF file into pages with text and images using PyMuPDF.
        Falls back to text-only extraction if PyMuPDF fails.
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            List of page dictionaries with text and image data
        """
        try:
            logger.info(f"Processing PDF: {pdf_path}")
            
            pages_data = []
            
            # Try PyMuPDF for full extraction (text + images)
            try:
                doc = fitz.open(str(pdf_path))
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    
                    # Extract text
                    text = page.get_text()
                    
                    # Render page to image
                    # Use matrix to set DPI (2.0 = 144 DPI, which is good quality)
                    zoom = self.pdf_dpi / 72  # PyMuPDF uses 72 DPI by default
                    mat = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat)
                    
                    # Convert pixmap to PIL Image
                    img_data = pix.tobytes("png")
                    page_image = Image.open(BytesIO(img_data))
                    
                    # Resize if needed
                    resized_image = self.resize_image(page_image)
                    image_base64 = self.image_to_base64(resized_image)
                    
                    page_data = {
                        "page_number": page_num + 1,
                        "text": text.strip(),
                        "image_base64": image_base64,
                        "image_width": resized_image.width,
                        "image_height": resized_image.height,
                        "source_file": str(pdf_path),
                        "content_type": "pdf_page"
                    }
                    
                    pages_data.append(page_data)
                
                doc.close()
                logger.info(f"Processed {len(pages_data)} pages from PDF with PyMuPDF")
                return pages_data
                
            except Exception as pymupdf_error:
                logger.warning(f"PyMuPDF processing failed: {pymupdf_error}")
                logger.info(f"Falling back to text-only extraction for {pdf_path.name}")
                return self._process_pdf_text_only(pdf_path)
            
        except Exception as e:
            logger.error(f"Error processing PDF {pdf_path}: {e}")
            return []
    
    def _process_pdf_text_only(self, pdf_path: Path) -> List[Dict[str, Any]]:
        """
        Process PDF file extracting text only (fallback when Poppler is not available).
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            List of page dictionaries with text data only
        """
        try:
            pages_data = []
            
            with open(pdf_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                
                for page_num, pdf_page in enumerate(pdf_reader.pages):
                    # Extract text
                    text = pdf_page.extract_text()
                    
                    if text.strip():  # Only add if there's actual text
                        page_data = {
                            "page_number": page_num + 1,
                            "text": text.strip(),
                            "source_file": str(pdf_path),
                            "content_type": "pdf_text",
                            "note": "Text-only extraction (Poppler not available)"
                        }
                        pages_data.append(page_data)
            
            logger.info(f"Processed {len(pages_data)} pages (text-only) from PDF")
            return pages_data
            
        except Exception as e:
            logger.error(f"Error in text-only PDF processing {pdf_path}: {e}")
            return []
    
    def process_pptx(self, pptx_path: Path) -> List[Dict[str, Any]]:
        """
        Process PowerPoint file into slides with text and notes.
        
        Args:
            pptx_path: Path to PowerPoint file
            
        Returns:
            List of slide dictionaries with text and notes
        """
        try:
            logger.info(f"Processing PowerPoint: {pptx_path}")
            
            slides_data = []
            prs = Presentation(str(pptx_path))
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract text from all shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                # Extract notes
                notes_text = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text
                
                # Combine text
                combined_text = "\\n\\n".join(slide_text)
                if notes_text:
                    combined_text += f"\\n\\nNotes: {notes_text}"
                
                slide_data = {
                    "slide_number": slide_num,
                    "text": combined_text.strip(),
                    "source_file": str(pptx_path),
                    "content_type": "pptx_slide"
                }
                
                slides_data.append(slide_data)
            
            logger.info(f"Processed {len(slides_data)} slides from PowerPoint")
            return slides_data
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint {pptx_path}: {e}")
            return []
    
    def process_image(self, image_path: Path) -> Optional[Dict[str, Any]]:
        """
        Process standalone image file.
        
        Args:
            image_path: Path to image file
            
        Returns:
            Image data dictionary or None if failed
        """
        try:
            logger.info(f"Processing image: {image_path}")
            
            with Image.open(image_path) as image:
                # Convert to RGB if necessary
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                resized_image = self.resize_image(image)
                image_base64 = self.image_to_base64(resized_image)
                
                image_data = {
                    "image_base64": image_base64,
                    "image_width": resized_image.width,
                    "image_height": resized_image.height,
                    "source_file": str(image_path),
                    "content_type": "image",
                    "filename": image_path.name
                }
                
                # Analyze image content using vision AI
                vision_data = self._analyze_image_content(image_path)
                image_data.update(vision_data)
                
                return image_data
                
        except Exception as e:
            logger.error(f"Error processing image {image_path}: {e}")
            return None
    
    def extract_html_content(self, html_content: str) -> Dict[str, Any]:
        """
        Extract text and image references from HTML content.
        
        Args:
            html_content: Raw HTML string
            
        Returns:
            Dictionary with extracted text and image URLs
        """
        try:
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract structured sections from details elements
            sections = self._extract_html_sections(soup)
            
            # Extract text (fallback for non-structured content)
            text = soup.get_text(separator=' ', strip=True)
            
            # Clean up whitespace
            text = re.sub(r'\s+', ' ', text).strip()
            
            # Extract image URLs
            image_urls = []
            seen_images: Dict[str, Dict[str, str]] = {}
            for img in soup.find_all('img'):
                src = img.get('src')
                if not src:
                    continue

                alt_text = img.get('alt', '')
                title_text = img.get('title', '')

                if src not in seen_images:
                    seen_images[src] = {
                        'src': src,
                        'alt': alt_text,
                        'title': title_text
                    }
                else:
                    existing = seen_images[src]

                    if alt_text and alt_text not in existing['alt']:
                        if existing['alt']:
                            existing['alt'] = f"{existing['alt']} | {alt_text}"
                        else:
                            existing['alt'] = alt_text

                    if title_text and title_text not in existing['title']:
                        if existing['title']:
                            existing['title'] = f"{existing['title']} | {title_text}"
                        else:
                            existing['title'] = title_text
            image_urls = list(seen_images.values())
            
            # Extract links
            links = []
            for link in soup.find_all('a'):
                href = link.get('href')
                if href:
                    links.append({
                        'href': href,
                        'text': link.get_text(strip=True)
                    })
            
            return {
                "text": text,
                "sections": sections,  # Add structured sections
                "image_urls": image_urls,
                "links": links,
                "content_type": "html"
            }
            
        except Exception as e:
            logger.error(f"Error extracting HTML content: {e}")
            return {"text": "", "sections": [], "image_urls": [], "links": [], "content_type": "html"}
    
    def _extract_html_sections(self, soup) -> List[Dict[str, str]]:
        """
        Extract sections from HTML details/summary structure.
        
        Args:
            soup: BeautifulSoup parsed HTML
            
        Returns:
            List of sections with headings and content
        """
        sections = []
        
        # Find top-level details elements (main sections)
        main_details = soup.find_all('details', recursive=False)
        
        # If no top-level details, look for any details elements
        if not main_details:
            main_details = soup.find_all('details')
        
        for detail in main_details:
            summary = detail.find('summary')
            if summary:
                # Extract heading from summary
                heading = summary.get_text(strip=True)
                # Clean up the heading text
                heading = re.sub(r'\s+', ' ', heading)
                
                # Extract content (everything in details except the summary)
                # Get all content excluding the summary
                content_parts = []
                for element in detail.children:
                    if element.name != 'summary':
                        if hasattr(element, 'get_text'):
                            content_parts.append(element.get_text(separator=' ', strip=True))
                        elif isinstance(element, str):
                            content_parts.append(element.strip())
                
                content = ' '.join(content_parts)
                content = re.sub(r'\s+', ' ', content).strip()
                
                if heading and content:
                    sections.append({
                        'heading': heading,
                        'content': content
                    })
        
        logger.info(f"Extracted {len(sections)} sections from HTML structure")
        return sections
    
    def chunk_text(self, text: str, metadata: Dict[str, Any] = None, html_sections: List[Dict[str, str]] = None) -> List[Dict[str, Any]]:
        """
        Split text into chunks with section-aware processing.
        
        Args:
            text: Text to chunk
            metadata: Additional metadata to include
            html_sections: Pre-extracted sections from HTML structure
            
        Returns:
            List of text chunks with metadata
        """
        if not text.strip():
            return []
        
        # Try section-aware chunking first
        sections = self._extract_sections(text, html_sections)
        if sections:
            return self._chunk_by_sections(sections, metadata)
        
        # Fall back to word-based chunking
        return self._chunk_by_words(text, metadata)
    
    def _extract_sections(self, text: str, html_sections: List[Dict[str, str]] = None) -> List[Dict[str, str]]:
        """
        Extract sections from text, preferring HTML structure if available.
        
        Args:
            text: Text to analyze
            html_sections: Pre-extracted sections from HTML structure
            
        Returns:
            List of sections with headings and content
        """
        # If we have HTML sections, use those directly
        if html_sections:
            logger.info(f"Using {len(html_sections)} sections from HTML structure")
            return html_sections
        
        # Fallback to text-based section detection for non-HTML content
        logger.info("Using fallback text-based section detection")
        
        # Known section headings from Canvas content (fallback only)
        known_headings = [
            "Why do we produce a 'Technical', 'Working', or' Construction' Drawing Pack?",
            "Who is responsible for the Technical Drawing Pack?",
            "When do we produce a Technical Drawing Pack?",
            "What is in an Architectural 'Technical', 'Working', or' Construction' Drawing Pack?",
            "Drawing Standards",
            "Construction 'Rules of Thumb'"
        ]
        
        # Alternative patterns for section detection
        question_patterns = [
            r'Why\s+do\s+we\s+produce.*?Drawing\s+Pack\?',
            r'Who\s+is\s+responsible.*?Drawing\s+Pack\?',
            r'When\s+do\s+we\s+produce.*?Drawing\s+Pack\?',
            r'What\s+is\s+in.*?Drawing\s+Pack\?',
            r'Drawing\s+Standards',
            r'Construction\s+[''"]\s*Rules\s+of\s+Thumb\s*[''""]'
        ]
        
        sections = []
        
        # Find all heading positions first
        heading_positions = []
        for heading in known_headings:
            pos = text.find(heading)
            if pos >= 0:
                heading_positions.append((pos, heading))
        
        # Sort by position
        heading_positions.sort()
        
        # Extract sections based on sorted positions
        for i, (pos, heading) in enumerate(heading_positions):
            section_start = pos
            
            # Find end of section (start of next section or end of text)
            if i + 1 < len(heading_positions):
                section_end = heading_positions[i + 1][0]
            else:
                section_end = len(text)
            
            # Extract section content
            section_text = text[section_start:section_end]
            
            # Split heading from content
            content_start = len(heading)
            content = section_text[content_start:].strip()
            
            # Remove duplicate heading if it appears again at start of content
            if content.startswith(heading):
                content = content[len(heading):].strip()
            
            sections.append({
                'heading': heading,
                'content': content
            })
        
        # If no known headings found, try pattern matching
        if not sections:
            for pattern in question_patterns:
                matches = list(re.finditer(pattern, text, re.IGNORECASE))
                for match in matches:
                    heading_start = match.start()
                    heading_end = match.end()
                    
                    # Find end of this section
                    section_end = len(text)
                    for other_pattern in question_patterns:
                        if other_pattern != pattern:
                            next_matches = list(re.finditer(other_pattern, text[heading_end:], re.IGNORECASE))
                            if next_matches:
                                section_end = heading_end + next_matches[0].start()
                                break
                    
                    heading = text[heading_start:heading_end]
                    content = text[heading_end:section_end].strip()
                    
                    sections.append({
                        'heading': heading,
                        'content': content
                    })
        
        # Sort sections by their position in the text
        if sections:
            sections.sort(key=lambda s: text.find(s['heading']))
            logger.info(f"Extracted {len(sections)} sections from text")
            return sections
        
        return []
    
    def _chunk_by_sections(self, sections: List[Dict[str, str]], metadata: Dict[str, Any] = None) -> List[Dict[str, Any]]:
        """
        Create chunks based on sections.
        
        Args:
            sections: List of sections with headings and content
            metadata: Additional metadata to include
            
        Returns:
            List of section-based chunks
        """
        chunks = []
        
        for i, section in enumerate(sections):
            heading = section['heading']
            content = section['content']
            
            # Create a chunk for the section heading itself
            heading_chunk = {
                "text": heading,
                "chunk_index": len(chunks),
                "section_index": i,
                "content_type": "section_heading",
                "is_section_heading": True
            }
            
            if metadata:
                heading_chunk.update(metadata)
            
            chunks.append(heading_chunk)
            
            # Create chunks for the section content
            if content:
                # If content is short enough, keep it as one chunk
                content_words = content.split()
                if len(content_words) <= self.chunk_size:
                    content_chunk = {
                        "text": f"{heading}\n\n{content}",
                        "chunk_index": len(chunks),
                        "section_index": i,
                        "section_heading": heading,
                        "content_type": "section_content"
                    }
                    
                    if metadata:
                        content_chunk.update(metadata)
                    
                    chunks.append(content_chunk)
                else:
                    # Split long content into sub-chunks but preserve section context
                    content_sub_chunks = self._chunk_by_words(content, metadata)
                    for j, sub_chunk in enumerate(content_sub_chunks):
                        sub_chunk.update({
                            "chunk_index": len(chunks),
                            "section_index": i,
                            "section_heading": heading,
                            "content_type": "section_content",
                            "sub_chunk_index": j,
                            "text": f"{heading}\n\n{sub_chunk['text']}"  # Include heading for context
                        })
                        chunks.append(sub_chunk)
        
        logger.info(f"Created {len(chunks)} section-aware chunks from {len(sections)} sections")
        return chunks
    
    def _chunk_by_words(self, text: str, metadata: Dict[str, Any] = None) -> List[Dict[str, Any]]:
        """
        Traditional word-based chunking (fallback method).
        
        Args:
            text: Text to chunk
            metadata: Additional metadata to include
            
        Returns:
            List of word-based chunks
        """
        chunks = []
        words = text.split()
        
        for i in range(0, len(words), self.chunk_size - self.chunk_overlap):
            chunk_words = words[i:i + self.chunk_size]
            chunk_text = ' '.join(chunk_words)
            
            chunk_data = {
                "text": chunk_text,
                "chunk_index": len(chunks),
                "start_word": i,
                "end_word": min(i + self.chunk_size, len(words)),
                "content_type": "text_chunk"
            }
            
            if metadata:
                chunk_data.update(metadata)
            
            chunks.append(chunk_data)
        
        return chunks
    
    def process_content_item(self, content_item: Dict[str, Any], parent_module: Optional[str] = None) -> List[Dict[str, Any]]:
        """
        Process a single content item (page, assignment, module, or file).
        
        Args:
            content_item: Content item from ingestion
            parent_module: Optional Canvas module title providing hierarchy context
            
        Returns:
            List of processed content segments
        """
        processed_segments = []
        
        try:
            # Determine effective parent module for this content and any nested segments
            parent_module_value = content_item.get("parent_module") or parent_module or ""
            
            if content_item["type"] == "page":
                # Process HTML page
                html_content = self.extract_html_content(content_item.get("body", ""))
                
                # Create text chunks
                text_chunks = self.chunk_text(
                    html_content["text"],
                    {
                        "source_id": content_item["id"],
                        "source_type": "page",
                        "title": content_item["title"],
                        "url": content_item["url"],
                        "created_at": content_item.get("created_at"),
                        "updated_at": content_item.get("updated_at"),
                        "parent_module": parent_module_value
                    },
                    html_content.get("sections", [])  # Pass HTML sections
                )
                processed_segments.extend(text_chunks)
                
                # Add image references with vision analysis
                for img_url in html_content["image_urls"]:
                    img_segment = {
                        "content_type": "image_reference",
                        "image_url": img_url["src"],
                        "alt_text": img_url["alt"],
                        "title": img_url["title"],
                        "source_id": content_item["id"],
                        "source_type": "page",
                        "page_title": content_item["title"],
                        "url": content_item["url"],
                        "parent_module": parent_module_value
                    }
                    
                    # Add vision analysis for the image URL
                    vision_data = self._analyze_image_content(img_url["src"], img_url["alt"])
                    img_segment.update(vision_data)
                    
                    processed_segments.append(img_segment)
            
            elif content_item["type"] == "assignment":
                # Process assignment (similar to page)
                html_content = self.extract_html_content(content_item.get("body", ""))
                
                # Create text chunks with assignment-specific metadata
                text_chunks = self.chunk_text(
                    html_content["text"],
                    {
                        "source_id": content_item["id"],
                        "source_type": "assignment",
                        "title": content_item["title"],
                        "url": content_item.get("html_url", ""),
                        "due_at": content_item.get("due_at"),
                        "points_possible": content_item.get("points_possible"),
                        "parent_module": parent_module_value
                    },
                    html_content.get("sections", [])
                )
                processed_segments.extend(text_chunks)
                
                # Process images in assignment
                for img_url in html_content["image_urls"]:
                    img_segment = {
                        "content_type": "image_reference",
                        "image_url": img_url["src"],
                        "alt_text": img_url["alt"],
                        "title": img_url["title"],
                        "source_id": content_item["id"],
                        "source_type": "assignment",
                        "page_title": content_item["title"],
                        "url": content_item.get("html_url", ""),
                        "parent_module": parent_module_value
                    }
                    vision_data = self._analyze_image_content(img_url["src"], img_url["alt"])
                    img_segment.update(vision_data)
                    processed_segments.append(img_segment)
            
            elif content_item["type"] == "module":
                # Process module structure
                module_text = f"Module: {content_item['title']}\n\n"
                
                # Add module items
                if "items" in content_item:
                    for item in content_item["items"]:
                        item_type = item.get("type", "unknown")
                        item_title = item.get("title", "Untitled")
                        module_text += f"- [{item_type}] {item_title}\n"
                        
                        # If item has embedded page content, process it
                        if "content" in item and isinstance(item["content"], dict):
                            # Recursively process the page content
                            page_segments = self.process_content_item(
                                item["content"],
                                parent_module=content_item["title"]
                            )
                            processed_segments.extend(page_segments)
                        
                        # If item is a File with downloaded content, process it
                        elif item.get("type") == "File" and "file_path" in item:
                            # Create a file content item and process it
                            file_content_item = {
                                "type": "file",
                                "path": item["file_path"],
                                "id": item["id"],
                                "title": item.get("title", "Untitled"),
                                "content_type": item.get("content_type", "application/octet-stream"),
                                "filename": item.get("filename", ""),
                                "parent_module": content_item["title"],
                                "url": item.get("html_url") or item.get("url") or item.get("file_url") or item.get("download_url", "")
                            }
                            file_segments = self.process_content_item(file_content_item)
                            processed_segments.extend(file_segments)
                
                # Create a single chunk for the module overview
                module_chunk = {
                    "text": module_text,
                    "content_type": "module_overview",
                    "source_id": content_item["id"],
                    "source_type": "module",
                    "title": content_item["title"],
                    "position": content_item.get("position", 0),
                    "parent_module": content_item["title"]
                }
                processed_segments.append(module_chunk)
            
            elif content_item["type"] == "file":
                file_path = Path(content_item["path"])
                
                if content_item["content_type"] == "application/pdf":
                    # Process PDF with PyMuPDF (text + images) or fallback to text-only
                    pdf_pages = self.process_pdf(file_path)
                    
                    for page_data in pdf_pages:
                        # Add metadata
                        page_data.update({
                            "source_id": content_item["id"],
                            "source_type": "file",
                            "filename": content_item["filename"],
                            "file_url": content_item.get("url", ""),
                            "parent_module": content_item.get("parent_module", ""),
                            "created_at": content_item.get("created_at"),
                            "updated_at": content_item.get("updated_at")
                        })
                        
                        # Create text chunks if there's text
                        if page_data["text"]:
                            text_chunks = self.chunk_text(
                                page_data["text"],
                                {
                                    "source_id": content_item["id"],
                                    "source_type": "pdf_text",
                                    "filename": content_item["filename"],
                                    "page_number": page_data["page_number"],
                                    "parent_module": content_item.get("parent_module", ""),
                                    "file_url": content_item.get("url", "")
                                }
                            )
                            processed_segments.extend(text_chunks)
                        
                        # Add the page as a multimodal segment
                        processed_segments.append(page_data)
                
                elif content_item["content_type"] in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", 
                                                       "application/vnd.ms-powerpoint"]:
                    # Process PowerPoint (.pptx or .ppt)
                    pptx_slides = self.process_pptx(file_path)
                    
                    for slide_data in pptx_slides:
                        # Add metadata
                        slide_data.update({
                            "source_id": content_item["id"],
                            "source_type": "file",
                            "filename": content_item["filename"],
                            "file_url": content_item.get("url", ""),
                            "parent_module": content_item.get("parent_module", ""),
                            "created_at": content_item.get("created_at"),
                            "updated_at": content_item.get("updated_at")
                        })
                        
                        # Create text chunks from slide content
                        if slide_data["text"]:
                            text_chunks = self.chunk_text(
                                slide_data["text"],
                                {
                                    "source_id": content_item["id"],
                                    "source_type": "pptx_text",
                                    "filename": content_item["filename"],
                                    "slide_number": slide_data["slide_number"],
                                    "parent_module": content_item.get("parent_module", ""),
                                    "file_url": content_item.get("url", "")
                                }
                            )
                            processed_segments.extend(text_chunks)
                        
                        # Add the slide data
                        processed_segments.append(slide_data)
                
                elif content_item["content_type"].startswith("image/"):
                    # Process image
                    image_data = self.process_image(file_path)
                    
                    if image_data:
                        image_data.update({
                            "source_id": content_item["id"],
                            "source_type": "file",
                            "file_url": content_item.get("url", ""),
                            "parent_module": content_item.get("parent_module", ""),
                            "created_at": content_item.get("created_at"),
                            "updated_at": content_item.get("updated_at")
                        })
                        processed_segments.append(image_data)
        
        except Exception as e:
            logger.error(f"Error processing content item {content_item.get('id', 'unknown')}: {e}")
        
        return processed_segments
    
    def process_course_content(self, metadata_path: Path) -> List[Dict[str, Any]]:
        """
        Process all content from a course metadata file.
        
        Supports multiple metadata formats:
        - Pages: {"pages": [...], "files": [...]} or {"page": {...}, "files": [...]}
        - Assignments: {"assignment": {...}, "files": [...]}
        - Modules: {"modules": [...]}
        
        Args:
            metadata_path: Path to course metadata JSON file
            
        Returns:
            List of all processed content segments
        """
        logger.info(f"Processing course content from: {metadata_path}")
        
        with open(metadata_path, 'r', encoding='utf-8') as f:
            course_metadata = json.load(f)
        
        all_segments = []
        
        # Handle different metadata formats
        if "pages" in course_metadata:
            # Course-wide format: {"pages": [...], "files": [...]}
            for page in course_metadata.get("pages", []):
                segments = self.process_content_item(page)
                all_segments.extend(segments)
                
        elif "page" in course_metadata:
            # Single-page format: {"page": {...}, "files": [...]}
            page = course_metadata["page"]
            segments = self.process_content_item(page)
            all_segments.extend(segments)
            
        elif "assignment" in course_metadata:
            # Assignment format: {"assignment": {...}, "files": [...]}
            assignment = course_metadata["assignment"]
            segments = self.process_content_item(assignment)
            all_segments.extend(segments)
            
        elif "modules" in course_metadata:
            # Modules format: {"modules": [...]}
            for module in course_metadata.get("modules", []):
                segments = self.process_content_item(module)
                all_segments.extend(segments)
        
        # Process files (common to pages and assignments)
        for file_item in course_metadata.get("files", []):
            segments = self.process_content_item(file_item)
            all_segments.extend(segments)
        
        # Save processed content
        output_path = settings.processed_data_dir / f"processed_{metadata_path.stem}.json"
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(all_segments, f, indent=2, ensure_ascii=False)
        
        logger.info(f"Processed {len(all_segments)} content segments")
        return all_segments
